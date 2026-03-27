"""
synco 투자 피치덱 생성 스크립트
python-pptx를 사용하여 15장 슬라이드의 사업계획서를 생성합니다.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 디자인 시스템 ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 딥 네이비
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)       # 카드 배경
ACCENT = RGBColor(0x00, 0xD4, 0xAA)        # 민트 그린 (synco 브랜드)
ACCENT2 = RGBColor(0x5B, 0x8D, 0xEF)       # 블루
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9A, 0xA0, 0xB4)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
RED_SOFT = RGBColor(0xFF, 0x6B, 0x6B)

SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="맑은 고딕", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    # line spacing
    from pptx.oxml.ns import qn
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=16,
                      default_color=WHITE, font_name="맑은 고딕", line_spacing=1.3):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)
    # Round corners
    shape.adjustments[0] = 0.05
    return shape


def add_divider(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_number(slide, num):
    add_text_box(slide, Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4),
                 str(num), font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)

add_text_box(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "synco", font_size=60, color=ACCENT, bold=True)
add_text_box(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             "AI가 기회를 발굴하고, 사람이 신뢰로 연결합니다",
             font_size=28, color=WHITE, bold=False)
add_divider(s, Inches(1.5), Inches(3.8), Inches(2))
add_text_box(s, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             "보험 설계사 28만 명의 CEO 네트워크를 AI 비즈니스 매칭으로 활성화하는 플랫폼",
             font_size=16, color=GRAY)
add_text_box(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
             "Pre-Seed 투자 제안서  |  2026", font_size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 2)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Executive Summary", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# 4 cards
cards = [
    ("문제", "보험 법인영업으로 구축된\nCEO 직통 DB가\n보험 계약 후 유휴 자산화", "💤"),
    ("솔루션", "무료 AI CRM으로 설계사 확보\nAI가 매칭 기회 발굴\n설계사가 대면으로 연결", "🔗"),
    ("시장", "SAM 800~2,900억원\nBEP 유료 CEO 17명\nYear 1 매출 ~1.6억원", "📊"),
    ("팀", "보험영업 8년 + 풀스택 7년\n공동창업자: GA 2개 계약\nCEO DB 4,000명 보유", "👥"),
]

for i, (title, body, icon) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x, Inches(1.5), Inches(2.8), Inches(4.5))
    add_text_box(s, x + Inches(0.3), Inches(1.7), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=ACCENT, bold=True)
    add_text_box(s, x + Inches(0.3), Inches(2.4), Inches(2.3), Inches(3.2),
                 body, font_size=15, color=LIGHT_GRAY, line_spacing=1.5)

add_text_box(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
             "Pre-seed 3억원  ·  18개월  ·  GA 5개 + CEO 5,000명 + CEO 과금 가설 검증",
             font_size=16, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: INSIGHT — 잠든 자산
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 3)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Unique Insight — 잠든 자산", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Main insight
add_rect(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.8), border=True)
add_text_box(s, Inches(1.2), Inches(1.6), Inches(10.5), Inches(1.5),
             "전국 GA 소속 설계사 288,000명이 법인영업으로 구축한 CEO 직통 전화번호 + 대면 관계 DB.\n"
             "GA가 수천만 원을 투자해서 만든 자산인데, 보험 계약이 끝나면 추가 수익 = 0원.",
             font_size=18, color=WHITE, line_spacing=1.5)

# Flow diagram using text
flow_items = [
    ("GA가 DB 구매", "수천만 원 투자"),
    ("FC가 법인영업", "CEO 대면 관계 구축"),
    ("보험 계약 완료", "추가 수익 0원"),
    ("DB 유휴 자산화", "가치가 잠듦"),
]

for i, (main, sub) in enumerate(flow_items):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x, Inches(3.7), Inches(2.6), Inches(1.5))
    add_text_box(s, x + Inches(0.2), Inches(3.85), Inches(2.2), Inches(0.5),
                 main, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(4.35), Inches(2.2), Inches(0.5),
                 sub, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, x + Inches(2.65), Inches(4.0), Inches(0.5), Inches(0.5),
                     "→", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Bottom message
add_rect(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.0), fill_color=RGBColor(0x15, 0x2E, 0x1E), border=False)
add_text_box(s, Inches(1.2), Inches(5.85), Inches(10.5), Inches(0.7),
             "synco는 이 잠든 자산을 깨웁니다. AI가 발굴하고, 사람이 성사시킵니다.",
             font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: PROBLEM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 4)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Problem", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Left: FC problems
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.2))
add_text_box(s, Inches(1.2), Inches(1.6), Inches(4.8), Inches(0.5),
             "설계사(FC)의 문제", font_size=20, color=ACCENT2, bold=True)

fc_problems = [
    "보험 재구매율 ≈ 0%  →  한번 팔면 관계 사장",
    "보험 외 줄 게 없어서  →  연락하면 차단",
    "5년간 쌓은 CEO 네트워크  →  보험 판매에만 소진",
    "관리 도구 = 카카오톡 메모 + 엑셀",
]
for j, prob in enumerate(fc_problems):
    add_text_box(s, Inches(1.2), Inches(2.3 + j * 0.9), Inches(4.8), Inches(0.8),
                 f"•  {prob}", font_size=14, color=LIGHT_GRAY, line_spacing=1.3)

# Right: CEO problems
add_rect(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.2))
add_text_box(s, Inches(7.2), Inches(1.6), Inches(4.8), Inches(0.5),
             "CEO(중소기업 대표)의 문제", font_size=20, color=ACCENT2, bold=True)

ceo_problems = [
    "설계사 전화 = 보험 팔려고  →  거부감",
    "유휴 설비, 재고, 유통채널  →  활용 채널 없음",
    "비즈니스 인맥 확장  →  100% 지인 소개 + 우연",
    "새 사업 기회 탐색  →  체계적 채널 부재",
]
for j, prob in enumerate(ceo_problems):
    add_text_box(s, Inches(7.2), Inches(2.3 + j * 0.9), Inches(4.8), Inches(0.8),
                 f"•  {prob}", font_size=14, color=LIGHT_GRAY, line_spacing=1.3)

# Bottom
add_text_box(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             "양쪽 모두 가치를 못 얻고 있습니다. FC에겐 관계 유지 명분이 없고, CEO에겐 기회 발견 채널이 없습니다.",
             font_size=16, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: SOLUTION — 양면 플랫폼
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 5)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Solution — 양면 플랫폼", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Left: FC CRM
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3.0))
add_text_box(s, Inches(1.2), Inches(1.55), Inches(4.5), Inches(0.5),
             "① 설계사용 무료 AI CRM", font_size=20, color=ACCENT, bold=True)
fc_sol = [
    "카카오톡 + 엑셀 대체",
    "AI가 기업 뉴스 수집 + 미팅 브리핑 생성",
    "GA의 단순 DB → '이 CEO에게 이 화두로 진입하세요' 변환",
    "무료 제공 = 트로이 목마",
]
for j, item in enumerate(fc_sol):
    add_text_box(s, Inches(1.2), Inches(2.3 + j * 0.55), Inches(4.8), Inches(0.5),
                 f"•  {item}", font_size=13, color=LIGHT_GRAY)

# Right: CEO App
add_rect(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(3.0))
add_text_box(s, Inches(7.2), Inches(1.55), Inches(4.5), Inches(0.5),
             "② CEO용 비즈니스 매칭 앱", font_size=20, color=ACCENT, bold=True)
ceo_sol = [
    "사업 정보 입력 → AI가 유휴 자원 분석",
    "맞춤 비즈니스 기회 자동 발굴",
    "관심 매칭만 크레딧으로 열람 (부분유료화)",
    "연결 = 이미 아는 설계사가 대면으로",
]
for j, item in enumerate(ceo_sol):
    add_text_box(s, Inches(7.2), Inches(2.3 + j * 0.55), Inches(4.8), Inches(0.5),
                 f"•  {item}", font_size=13, color=LIGHT_GRAY)

# Center highlight
add_rect(s, Inches(3.0), Inches(4.8), Inches(7.3), Inches(1.5), fill_color=RGBColor(0x15, 0x2E, 0x1E), border=True)
add_text_box(s, Inches(3.3), Inches(4.95), Inches(6.8), Inches(0.5),
             "AI가 발굴하고, 사람이 성사시킵니다", font_size=24, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(3.3), Inches(5.5), Inches(6.8), Inches(0.5),
             "콜드 메시지가 아닙니다. 이미 신뢰 관계가 있는 설계사의 웜 채널입니다.",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: 구체적 예시 (스토리텔링)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 6)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "How It Works — 실제 시나리오", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Story cards
story_steps = [
    ("화성 제조업체", "매출 85억, 신공장 증설\n가동률 60%\n40% 유휴 라인", ACCENT2),
    ("수원 소비재 회사", "매출 50억\n위탁생산처 필요\n마땅한 곳 없음", ACCENT2),
    ("같은 설계사의 고객", "synco AI가\n매칭 기회 발견\n'유휴 라인 ↔ 위탁 수요'", ACCENT),
    ("설계사가 대면 연결", "'김 사장님, 좋은 분\n소개해드릴까요?'\n→ 수억 원 거래", YELLOW),
]

for i, (title, body, color) in enumerate(story_steps):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(1.5), Inches(2.9), Inches(3.5))
    # Step number
    add_text_box(s, x + Inches(0.2), Inches(1.6), Inches(0.5), Inches(0.4),
                 str(i+1), font_size=28, color=color, bold=True)
    add_text_box(s, x + Inches(0.2), Inches(2.1), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(s, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(2.0),
                 body, font_size=13, color=LIGHT_GRAY, line_spacing=1.5)
    if i < 3:
        add_text_box(s, x + Inches(2.9), Inches(2.7), Inches(0.4), Inches(0.5),
                     "→", font_size=24, color=ACCENT)

add_text_box(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
             "서로의 존재를 몰랐던 두 CEO를, 이미 신뢰 관계가 있는 설계사가 연결합니다.\n"
             "이 한 마디가 수억 원의 거래가 됩니다.",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: TRACTION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 7)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Traction — 검증 현황", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Traction items
traction = [
    ("2개 GA 계약 체결", "공동창업자가 현업 영업자로서 직접 파트너십 체결. Day 1 실행 가능.", ACCENT),
    ("CEO DB 3,000~4,000명", "공동창업자 본인 DB로 즉시 프로토타입 테스트 가능.", ACCENT),
    ("CEO 대면 반응 확인", "보험 제안은 거부하지만, 비즈니스 매칭은 긍정적.", ACCENT2),
    ("FC 인터뷰 확인", "쓸만한 CRM이 없다. 기존 도구(카톡/엑셀)에 불만 크다.", ACCENT2),
]

for i, (title, desc, color) in enumerate(traction):
    y = Inches(1.4 + i * 1.2)
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.0))
    # Checkmark
    add_text_box(s, Inches(1.2), y + Inches(0.05), Inches(0.5), Inches(0.5),
                 "✓", font_size=24, color=color, bold=True)
    add_text_box(s, Inches(1.8), y + Inches(0.1), Inches(4), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(s, Inches(1.8), y + Inches(0.55), Inches(10), Inches(0.4),
                 desc, font_size=13, color=GRAY)

# BNI benchmark
add_rect(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), fill_color=RGBColor(0x15, 0x2E, 0x1E))
add_text_box(s, Inches(1.2), Inches(6.1), Inches(10.5), Inches(0.7),
             "벤치마크: BNI 한국 2,300명이 리퍼럴로 연간 수천억 거래 → FC 30만 명이 같은 역할을 하면?",
             font_size=15, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: MARKET SIZE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 8)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Market Size", font_size=28, color=ACCENT, bold=True)
add_text_box(s, Inches(6), Inches(0.4), Inches(6), Inches(0.6),
             "Bottom-up 계산", font_size=14, color=GRAY, alignment=PP_ALIGN.RIGHT)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# TAM / SAM / SOM
market_data = [
    ("TAM", "6,000억 ~ 4.8조원", "도달 가능 Unique CEO 50만~100만명\nCEO 1인당 연간 120만~480만원 (구독+크레딧+딜 중개)", Inches(3.8)),
    ("SAM", "800억 ~ 2,900억원", "법인영업 FC ~10만 명이\n접근 가능한 CEO 기반", Inches(3.0)),
    ("SOM", "~1.6억원 (Year 1)", "GA 5개 × CEO 20,000명\n유료 전환 3%", Inches(2.2)),
]

for i, (label, amount, desc, bar_w) in enumerate(market_data):
    y = Inches(1.4 + i * 1.8)
    # Label
    add_text_box(s, Inches(0.8), y, Inches(1.2), Inches(0.5),
                 label, font_size=22, color=ACCENT, bold=True)
    # Amount
    add_text_box(s, Inches(2.2), y, Inches(4), Inches(0.5),
                 amount, font_size=22, color=WHITE, bold=True)
    # Bar
    bar_color = ACCENT if i == 0 else ACCENT2 if i == 1 else YELLOW
    add_rect(s, Inches(2.2), y + Inches(0.5), bar_w, Inches(0.15), fill_color=bar_color)
    # Description
    add_text_box(s, Inches(7), y, Inches(5.5), Inches(1.2),
                 desc, font_size=13, color=GRAY, line_spacing=1.4)

# BEP highlight
add_rect(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), border=True)
add_text_box(s, Inches(1.2), Inches(6.1), Inches(10.5), Inches(0.7),
             "손익분기점 = 유료 CEO 17명  (월 비용 205만원 ÷ 공헌이익 12.65만원)",
             font_size=18, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: BUSINESS MODEL
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 9)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Business Model — 3-Layer", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

layers = [
    ("Layer 1", "트로이 목마", "0~6개월",
     ["FC 무료 AI CRM 제공", "CEO 무료 CRM(인맥/고객 관리) 제공", "양쪽 가입 유도 = 네트워크 구축"],
     ACCENT2),
    ("Layer 2", "과금", "6~18개월",
     ["CEO 구독: 월 9.9만원(Premium) / 39.9만원(Enterprise)",
      "크레딧: 매칭 상세 열람 건당 5,000~25,000원",
      "딜 중개 수수료: 거래액의 3~10% ← 가장 큰 시장",
      "FC 구독: Pro 3.9만원 / Premium 9.9만원"],
     ACCENT),
    ("Layer 3", "금융상품 유통", "18개월+",
     ["CEO DB 10만명+ = 고자산 고객 풀",
      "투자자문, 금융상품을 FC가 대면 유통",
      "로보어드바이저 = 첫 번째 상품"],
     YELLOW),
]

for i, (name, subtitle, period, items, color) in enumerate(layers):
    x = Inches(0.8 + i * 4.1)
    add_rect(s, x, Inches(1.3), Inches(3.8), Inches(4.8))
    add_text_box(s, x + Inches(0.3), Inches(1.45), Inches(3.2), Inches(0.4),
                 f"{name}: {subtitle}", font_size=18, color=color, bold=True)
    add_text_box(s, x + Inches(0.3), Inches(1.9), Inches(3.2), Inches(0.3),
                 period, font_size=12, color=GRAY)
    add_divider(s, x + Inches(0.3), Inches(2.25), Inches(3.2), color=color)
    for j, item in enumerate(items):
        add_text_box(s, x + Inches(0.3), Inches(2.5 + j * 0.7), Inches(3.3), Inches(0.65),
                     f"•  {item}", font_size=12, color=LIGHT_GRAY, line_spacing=1.3)

# Bottom: key message
add_text_box(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
             "핵심: 매칭은 DB를 활성화시키는 도구. 활성화된 DB 위에 무한한 비즈니스가 파생됩니다.",
             font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: COMPETITIVE LANDSCAPE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 10)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Competitive Landscape", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Table header
headers = ["", "매칭 방식", "성사 방식", "한계", "synco 차별화"]
header_widths = [Inches(1.8), Inches(2.0), Inches(2.0), Inches(2.8), Inches(2.8)]

x_pos = Inches(0.5)
for col, (header, w) in enumerate(zip(headers, header_widths)):
    add_rect(s, x_pos, Inches(1.3), w, Inches(0.5), fill_color=RGBColor(0x25, 0x35, 0x55))
    add_text_box(s, x_pos + Inches(0.1), Inches(1.33), w - Inches(0.2), Inches(0.45),
                 header, font_size=12, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    x_pos += w

competitors = [
    ("리멤버\n250만 회원", "AI 추천", "콜드 메시지", "신뢰 부재\n대면 없음", "FC 웜 채널\n신뢰 내재"),
    ("BNI\n2,300명", "사람 소개", "대면 리퍼럴", "확장 불가\n주 1회 모임", "AI 확장성\n매일 매칭"),
    ("링크드인", "프로필 매칭", "콜드 InMail", "한국 CEO\n사용률 낮음", "FC 대면\n한국 특화"),
    ("지인 소개", "우연", "우연", "체계 없음\n재현 불가", "AI 체계화\n데이터 기반"),
]

for row, (name, match, close, limit, diff) in enumerate(competitors):
    y = Inches(1.85 + row * 1.1)
    bg_c = BG_CARD if row % 2 == 0 else RGBColor(0x12, 0x1D, 0x32)
    x_pos = Inches(0.5)
    for col, (text, w) in enumerate(zip([name, match, close, limit, diff], header_widths)):
        add_rect(s, x_pos, y, w, Inches(1.0), fill_color=bg_c)
        clr = WHITE if col == 0 else LIGHT_GRAY if col < 4 else ACCENT
        bld = col == 0 or col == 4
        add_text_box(s, x_pos + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.9),
                     text, font_size=11, color=clr, bold=bld, alignment=PP_ALIGN.CENTER,
                     line_spacing=1.3)
        x_pos += w

add_text_box(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             "리멤버가 시장을 교육해주고 있지만, 콜드 매칭입니다. synco는 웜 채널입니다. 같은 시장, 다른 접근.",
             font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: WHY NOW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 11)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Why Now", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

why_now = [
    ("AI CRM 폭발적 성장", "글로벌 AI CRM 시장 CAGR 36%\n14.9조 → 69.8조원\n2025~2027년이 초기→주류 전환 결정적 시기",
     "📈", ACCENT),
    ("GA의 구조적 변화", "GA협회 2035년 '금융판매전문회사' 비전 공식 채택\n보험 판매 → 종합 금융/비즈니스 연결\nsynco는 이 비전에 정확히 정렬",
     "🏢", ACCENT2),
    ("리멤버 매칭 시장 진입", "250만 회원의 리멤버가 2026.3 매칭 런칭\n시장을 교육해주지만 = 콜드 매칭\nsynco = 웜 매칭. 같은 시장, 다른 접근",
     "🔥", YELLOW),
]

for i, (title, desc, icon, color) in enumerate(why_now):
    y = Inches(1.4 + i * 1.7)
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.5))
    add_text_box(s, Inches(1.3), y + Inches(0.15), Inches(4), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(s, Inches(1.3), y + Inches(0.65), Inches(10.5), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, line_spacing=1.4)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: TEAM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 12)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Team", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Founder
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.5))
add_text_box(s, Inches(1.3), Inches(1.6), Inches(4.5), Inches(0.5),
             "창업자 (CEO & CTO)", font_size=20, color=ACCENT, bold=True)

founder_items = [
    ("보험영업 8년", "FC의 습성과 고충을 체감적으로 이해"),
    ("풀스택 개발 7년", "투자/자산관리 프로그램 직접 개발"),
    ("로보어드바이저 운용 중", "Layer 3의 첫 번째 상품이 이미 존재"),
    ("기획 + 개발 1인", "MVP를 외주 없이 빠르게 구축 가능"),
]
for j, (title, desc) in enumerate(founder_items):
    add_text_box(s, Inches(1.3), Inches(2.3 + j * 0.85), Inches(4.5), Inches(0.4),
                 f"▸ {title}", font_size=15, color=WHITE, bold=True)
    add_text_box(s, Inches(1.6), Inches(2.7 + j * 0.85), Inches(4.2), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# Co-founder
add_rect(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(4.5))
add_text_box(s, Inches(7.3), Inches(1.6), Inches(4.5), Inches(0.5),
             "공동창업자 (COO & Sales)", font_size=20, color=ACCENT, bold=True)

cofounder_items = [
    ("현업 금융 영업자", "보험/금융 현장에서 활동 중"),
    ("2개 GA 계약 체결", "Day 1 파트너십 실행 가능"),
    ("CEO DB 3,000~4,000명", "즉시 프로토타입 테스트 가능"),
    ("영업 네트워크", "추가 GA 확장 파이프라인 보유"),
]
for j, (title, desc) in enumerate(cofounder_items):
    add_text_box(s, Inches(7.3), Inches(2.3 + j * 0.85), Inches(4.5), Inches(0.4),
                 f"▸ {title}", font_size=15, color=WHITE, bold=True)
    add_text_box(s, Inches(7.6), Inches(2.7 + j * 0.85), Inches(4.2), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# Bottom
add_rect(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), fill_color=RGBColor(0x15, 0x2E, 0x1E))
add_text_box(s, Inches(1.2), Inches(6.3), Inches(10.5), Inches(0.6),
             "이 팀은 돈으로 살 수 없습니다. 8년간 보험 현장을 경험한 사람 + 지금 당장 GA에 가서 계약할 수 있는 사람.",
             font_size=15, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: FINANCIALS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 13)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Financial Projections", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Key metrics row
metrics = [
    ("BEP", "유료 CEO 17명", YELLOW),
    ("Year 1", "~1.6억원", WHITE),
    ("Year 2", "10~18억원", WHITE),
    ("Year 3", "30~55억원", WHITE),
    ("Gross Margin", "~83%", ACCENT),
]

for i, (label, value, color) in enumerate(metrics):
    x = Inches(0.5 + i * 2.5)
    add_rect(s, x, Inches(1.3), Inches(2.2), Inches(1.5))
    add_text_box(s, x + Inches(0.2), Inches(1.45), Inches(1.8), Inches(0.4),
                 label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(1.9), Inches(1.8), Inches(0.5),
                 value, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Revenue breakdown
add_text_box(s, Inches(0.8), Inches(3.2), Inches(5), Inches(0.5),
             "Year 1 수익 구조", font_size=18, color=WHITE, bold=True)

rev_items = [
    ("CEO 구독 + 크레딧", "~1.3억원", "81%"),
    ("FC 구독", "~2,070만원", "13%"),
    ("딜 중개 수수료", "~670만원", "4%"),
    ("합계", "~1.6억원", "100%"),
]

for i, (item, amount, pct) in enumerate(rev_items):
    y = Inches(3.8 + i * 0.55)
    clr = YELLOW if i == 3 else LIGHT_GRAY
    bld = i == 3
    add_text_box(s, Inches(1.2), y, Inches(3.5), Inches(0.45),
                 item, font_size=14, color=clr, bold=bld)
    add_text_box(s, Inches(4.8), y, Inches(2), Inches(0.45),
                 amount, font_size=14, color=clr, bold=bld, alignment=PP_ALIGN.RIGHT)
    add_text_box(s, Inches(7.0), y, Inches(1), Inches(0.45),
                 pct, font_size=14, color=GRAY, alignment=PP_ALIGN.RIGHT)

# Cost structure
add_text_box(s, Inches(8.5), Inches(3.2), Inches(4), Inches(0.5),
             "비용 구조 (Phase별 월 비용)", font_size=18, color=WHITE, bold=True)

cost_items = [
    ("Phase 1 (0~6개월)", "월 45만원"),
    ("Phase 2 (6~12개월)", "월 185만원"),
    ("Phase 3 (12~18개월)", "월 755만원"),
]

for i, (phase, cost) in enumerate(cost_items):
    y = Inches(3.8 + i * 0.55)
    add_text_box(s, Inches(8.8), y, Inches(2.5), Inches(0.45),
                 phase, font_size=13, color=LIGHT_GRAY)
    add_text_box(s, Inches(11.3), y, Inches(1.5), Inches(0.45),
                 cost, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 14: THE ASK
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 14)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "The Ask", font_size=28, color=ACCENT, bold=True)
add_divider(s, Inches(0.8), Inches(1.0), Inches(11.5))

# Main ask
add_rect(s, Inches(3.5), Inches(1.3), Inches(6.3), Inches(1.2), border=True)
add_text_box(s, Inches(3.8), Inches(1.4), Inches(5.7), Inches(1.0),
             "Pre-seed 3억원  ·  18개월", font_size=32, color=YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

# Milestones
milestones = [
    ("Month 3", "MVP 출시 + GA 2개 + FC 100명 + CEO 500명", "플랫폼 작동 검증"),
    ("Month 6", "GA 5개 + CEO 5,000명 + 첫 유료 전환", "CEO WTP 가설 검증 (핵심)"),
    ("Month 12", "GA 20개 + CEO 30,000명 + MRR 1,000만원", "성장 모멘텀 확보"),
    ("Month 18", "CEO 50,000명 + Seed 라운드 준비", "스케일링 단계 진입"),
]

for i, (month, target, meaning) in enumerate(milestones):
    y = Inches(2.8 + i * 0.9)
    color_m = YELLOW if i == 1 else ACCENT
    add_rect(s, Inches(0.8), y, Inches(1.8), Inches(0.7), fill_color=RGBColor(0x25, 0x35, 0x55))
    add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(1.6), Inches(0.5),
                 month, font_size=16, color=color_m, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.8), y + Inches(0.1), Inches(5.5), Inches(0.5),
                 target, font_size=14, color=WHITE)
    add_text_box(s, Inches(8.5), y + Inches(0.1), Inches(4), Inches(0.5),
                 meaning, font_size=13, color=GRAY)

# Fund allocation
add_text_box(s, Inches(0.8), Inches(6.0), Inches(5), Inches(0.4),
             "자금 용도", font_size=16, color=WHITE, bold=True)

alloc = [("개발 인프라 + AI API", "40%"), ("GA 파트너십 + 운영", "30%"),
         ("마케팅 + CEO 획득", "20%"), ("법률/컴플라이언스", "10%")]

for i, (item, pct) in enumerate(alloc):
    x = Inches(0.8 + i * 3.1)
    add_text_box(s, x, Inches(6.4), Inches(2.8), Inches(0.4),
                 f"{item}: {pct}", font_size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: CLOSING
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text_box(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "synco", font_size=60, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1.5), Inches(3.0), Inches(10), Inches(1.5),
             "AI 시대에 비즈니스 모델은 카피할 수 있습니다.\n"
             "하지만 8년간 현장에서 체험한 가치관은 카피할 수 없습니다.\n"
             "사람은 유일한 존재이기 때문입니다.",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_divider(s, Inches(5.5), Inches(4.8), Inches(2.3))

add_text_box(s, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
             "AI가 기회를 발굴하고, 사람이 신뢰로 연결합니다.\n"
             "이것이 synco의 철학이고, 이 철학이 곧 해자입니다.",
             font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_text_box(s, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "감사합니다.", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = r"C:\Users\chaconne\Desktop\business\synco\07-pitch\synco-pitch-deck.pptx"
prs.save(output_path)
print(f"✅ 피치덱 생성 완료: {output_path}")
print(f"   슬라이드 수: {len(prs.slides)}")
